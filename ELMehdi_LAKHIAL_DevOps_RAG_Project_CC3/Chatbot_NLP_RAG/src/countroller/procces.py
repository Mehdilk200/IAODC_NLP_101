from .project import ProjectController
from .base import BaseController
import os
from models.enums.roles import roles
import re
import json
from pptx import Presentation
import pandas as pd
from pypdf import PdfReader
from langchain_core.documents import Document

from models.enums.extenctionEnum import extenction
from langchain_text_splitters import RecursiveCharacterTextSplitter

class ProcessController(BaseController):
        def __init__(self, project_id: roles):
                super().__init__()
                self.project_id = project_id.value if hasattr(project_id, 'value') else project_id
                self.project_file_path = ProjectController().get_project_path(project_id=self.project_id)
                self.code_patterns = [
                        r'def\s+\w+\(.*\):', 
                        r'class\s+\w+[:\(]', 
                        r'import\s+\w+',     
                        r'\{.*\}',           
                        r'console\.log',     
                        r'return\s+'         
                        ]
                self.combined_pattern = "(" + "|".join(self.code_patterns) + ")"
        
        def _load_txt(self, path: str):
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()
                return [Document(page_content=text, metadata={"source": os.path.basename(path), "path": path})]

        def _load_pdf(self, path: str):
                reader = PdfReader(path)
                docs = []
                for i, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        docs.append(Document(
                                page_content=text,
                                metadata={"source": os.path.basename(path), "path": path, "page": i + 1}
                        ))
                return docs

        def _load_docx(self, path: str):
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text)
                return [Document(page_content=text, metadata={"source": os.path.basename(path), "path": path})]
        def _load_excel(self, path: str):
                df = pd.read_excel(path)

                text = df.to_string(index=False)

                return [
                        Document(
                        page_content=text,
                        metadata={
                                "source": os.path.basename(path),
                                "path": path
                        }
                        )
                ]

        def _load_csv(self, path: str):
                df = pd.read_csv(path)
                text = df.to_string()
                return [Document(page_content=text, metadata={"source": os.path.basename(path), "path": path})]

        def _load_pptx(self, path: str):
                prs = Presentation(path)
                text = []

                for slide in prs.slides:
                        for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                        text.append(shape.text)

                return [Document(page_content="\n".join(text),
                                metadata={"source": os.path.basename(path), "path": path})]

        def _load_json(self, path: str):
                with open(path, "r", encoding="utf-8") as f:
                        data = json.load(f)

                text = json.dumps(data, indent=2)

                return [Document(page_content=text,
                                metadata={"source": os.path.basename(path), "path": path})]
                        
        def get_extension(self, file_name : str):
                return os.path.splitext(file_name)[-1] 

        def get_load_file(self, file_name: str):
                file_extension = self.get_extension(file_name=file_name)
                file_location  = os.path.join(self.project_file_path, file_name)

                loaders = {
                        extenction.TXT.value:  lambda: self._load_txt(file_location),
                        extenction.MD.value:   lambda: self._load_txt(file_location),
                        extenction.XML.value:  lambda: self._load_txt(file_location),
                        extenction.PDF.value:  lambda: self._load_pdf(file_location),
                        extenction.DOCX.value: lambda: self._load_docx(file_location),
                        extenction.XLS.value:  lambda: self._load_excel(file_location),
                        extenction.XLSX.value: lambda: self._load_excel(file_location),
                        extenction.CSV.value:  lambda: self._load_csv(file_location),
                        extenction.PPTX.value: lambda: self._load_pptx(file_location),
                        extenction.JSON.value: lambda: self._load_json(file_location),
                }

                loader = loaders.get(file_extension)
                if not loader:
                        raise ValueError(f"Unsupported: {file_extension}")

                return loader()


        def get_content_file(self, file_name: str):
                try:
                        content = self.get_load_file(file_name=file_name)
                        return content
                except Exception as e:
                        raise ValueError(f"Error loading file: {str(e)}")
        
        def process_file(self, file_name: str, chunk_size: int = 200, overlap_size: int = 20, file_content: list = None):
                if file_content is None:
                        file_content = []

                text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=chunk_size, 
                        chunk_overlap=overlap_size, 
                        length_function=len
                )
                
                file_content_text = [doc.page_content for doc in file_content]
                file_content_meta = [doc.metadata for doc in file_content]

                
                chunks = text_splitter.create_documents(file_content_text, metadatas=file_content_meta)

                return chunks

        def has_code(self, text: str) -> bool:
                if text is None:
                     return False
                if hasattr(text, 'page_content'):
                    text = text.page_content
                text_to_search = str(text)
                return bool(re.search(self.combined_pattern, text_to_search))

                
                